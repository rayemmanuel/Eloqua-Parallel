import fitz
from docx import Document
from pptx import Presentation
from google import genai
from google.genai import types
from config import GEMINI_API_KEY
from pydantic import BaseModel
import json, re

client = genai.Client(api_key=GEMINI_API_KEY)

def extract_text(file_path: str, file_type: str) -> str:
    if file_type == "pdf":
        doc = fitz.open(file_path)
        return "\n".join([page.get_text() for page in doc])
    elif file_type == "docx":
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    elif file_type == "pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    return ""


# ── Pydantic models for structured Gemini outputs ──────────────────────────────

class TalkingPointsResult(BaseModel):
    title: str
    talking_points: list[str]


class RelevanceResult(BaseModel):
    relevance_score: float
    relevance_feedback: str


class CoveragePoint(BaseModel):
    talking_point: str
    covered: bool
    confidence: str
    feedback: str


class CoverageResult(BaseModel):
    coverage_report: list[CoveragePoint]
    coverage_score: float
    coverage_feedback: str


# ── API Logic ──────────────────────────────────────────────────────────────────

def generate_talking_points(text: str) -> dict:
    prompt = f"""
You are an academic speech coach. Based on the following document content,
generate a structured list of 5 to 7 key talking points the student should
be able to explain verbally during a presentation.

Each talking point should be:
- Clear and concise (one sentence)
- Actionable (something the student can actually say out loud)
- Based strictly on the document content

Document content:
{text[:3000]}
"""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=TalkingPointsResult
            )
        )
        data = json.loads(response.text)
        return {
            "title": str(data.get("title", "Document Summary")),
            "talking_points": [str(tp) for tp in data.get("talking_points", [])]
        }
    except Exception as e:
        print(f"[DOCUMENT] Error generating talking points: {e}")
        return {
            "title": "Document Analysis",
            "talking_points": [
                "Explain the main objective of the presentation.",
                "Detail the key methodology or structure discussed.",
                "Summarize the main results or conclusions."
            ]
        }

def check_relevance(topic: str, transcript: str) -> dict:
    """
    Checks if the student's speech is relevant to the given topic.
    Used for Spontaneous Mode sessions.
    """
    if not topic.strip() or not transcript.strip():
        return {
            "relevance_score": 0,
            "relevance_feedback": "No topic or transcript available to evaluate."
        }

    prompt = f"""
You are an academic speech coach evaluating a student's impromptu speech.

The student was given this topic to speak about:
"{topic}"

This is the transcript of what the student actually said:
"{transcript}"

Evaluate how relevant and on-topic the student's speech was.
Consider:
- Did they actually address the topic?
- Did they stay on topic throughout?
- Did they provide a clear introduction, definition, example, or analysis related to the topic?
"""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=RelevanceResult
            )
        )
        data = json.loads(response.text)
        return {
            "relevance_score": float(data.get("relevance_score", 0.0)),
            "relevance_feedback": str(data.get("relevance_feedback", ""))
        }
    except Exception as e:
        print(f"[DOCUMENT] Error checking relevance: {e}")
        return {
            "relevance_score": 0,
            "relevance_feedback": "Could not verify topic relevance."
        }

def check_coverage(talking_points: list, transcript: str) -> dict:
    """
    Compares the user's speech transcript against the generated
    talking points and returns a per-point coverage report.
    Only called during Preparation Mode sessions.
    """
    if not talking_points or not transcript.strip():
        return {
            "coverage_report": [],
            "coverage_score": 0,
            "coverage_feedback": "No talking points or transcript available to evaluate."
        }

    prompt = f"""
You are an academic speech coach evaluating a student's presentation.

The student was given these talking points to cover:
{json.dumps(talking_points, indent=2)}

This is the transcript of what the student actually said:
"{transcript}"

For each talking point, determine if the student covered it in their speech.
Be fair — the student does not need to use exact words, just convey the idea clearly.
"""
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=CoverageResult
            )
        )
        data = json.loads(response.text)
        return {
            "coverage_report": [
                {
                    "talking_point": str(cp.get("talking_point", "")),
                    "covered": bool(cp.get("covered", False)),
                    "confidence": str(cp.get("confidence", "low")),
                    "feedback": str(cp.get("feedback", ""))
                }
                for cp in data.get("coverage_report", [])
            ],
            "coverage_score": float(data.get("coverage_score", 0.0)),
            "coverage_feedback": str(data.get("coverage_feedback", ""))
        }
    except Exception as e:
        print(f"[DOCUMENT] Error checking coverage: {e}")
        return {
            "coverage_report": [],
            "coverage_score": 0,
            "coverage_feedback": "Could not verify talking point coverage."
        }


def chat_with_coach(history: list, system_prompt: str) -> str:
    contents = []
    for msg in history:
        role = "user" if msg["role"] == "user" else "model"
        contents.append(types.Content(
            role=role,
            parts=[types.Part(text=msg["content"])]
        ))
    
    response = client.models.generate_content(
        model="gemini-2.5-flash",
        contents=contents,
        config=types.GenerateContentConfig(
            system_instruction=system_prompt,
            max_output_tokens=500,
        )
    )
    return response.text.strip()